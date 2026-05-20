
from __future__ import annotations

import re
from pathlib import Path
from typing import Optional

# Regex de tokenización (reutilizada de store.py)
_TOKEN_RE = re.compile(r"[\wáéíóúñüÁÉÍÓÚÑÜ]+", re.UNICODE)

# Extensiones de texto plano que se leen directamente
_PLAIN_TEXT_EXTENSIONS = {
    "txt", "md", "markdown", "rst", "csv", "json", "yaml", "yml",
    "xml", "toml", "ini", "cfg", "conf", "env", "log",
    "html", "htm", "css", "scss", "sass",
    "js", "jsx", "ts", "tsx", "vue", "svelte",
    "py", "pyw", "java", "kt", "c", "h", "cpp", "cc", "cs",
    "go", "rs", "rb", "php", "swift", "dart", "r", "sql",
    "sh", "bash", "zsh", "ps1", "bat", "lua", "pl",
    "scala", "ex", "exs", "hs", "clj",
}

# Número máximo de bytes a leer de un archivo de texto plano
_MAX_PLAIN_BYTES = 1 * 1024 * 1024  # 1 MB


def extraer_texto(ruta: str | Path) -> Optional[str]:
    """
    Extrae el contenido de texto de un archivo según su extensión.
    Retorna None si el tipo no es soportado o si ocurre un error.
    """
    ruta = Path(ruta)
    if not ruta.exists():
        return None

    extension = ruta.suffix.lower().lstrip(".")

    try:
        if extension == "pdf":
            return _extraer_pdf(ruta)
        elif extension in ("docx", "doc", "odt"):
            return _extraer_docx(ruta)
        elif extension in ("xlsx", "xls", "ods"):
            return _extraer_xlsx(ruta)
        elif extension in ("pptx", "ppt"):
            return _extraer_pptx(ruta)
        elif extension in ("png", "jpg", "jpeg", "gif", "webp", "bmp", "tiff"):
            return _extraer_imagen(ruta)
        elif extension in ("mp3", "wav", "ogg", "flac", "m4a", "aac"):
            return _extraer_audio(ruta)
        elif extension in ("mp4", "avi", "mkv", "mov", "webm"):
            return _extraer_video(ruta)
        elif extension in _PLAIN_TEXT_EXTENSIONS:
            return _extraer_texto_plano(ruta)
    except Exception as e:
        # Si falla la extracción, no bloqueamos la subida del archivo
        print(f"Error en extraer_texto para {ruta}: {e}")
        return None

    return None


def _extraer_pdf(ruta: Path) -> Optional[str]:
    """Extrae texto de un PDF usando pdfplumber."""
    try:
        import pdfplumber
        partes: list[str] = []
        with pdfplumber.open(str(ruta)) as pdf:
            for pagina in pdf.pages:
                texto = pagina.extract_text()
                if texto:
                    partes.append(texto.strip())
        return "\n".join(partes) if partes else None
    except ImportError:
        return None


def _extraer_docx(ruta: Path) -> Optional[str]:
    """Extrae texto de un archivo Word (.docx)."""
    try:
        import docx
        doc = docx.Document(str(ruta))
        partes = [parrafo.text for parrafo in doc.paragraphs if parrafo.text.strip()]
        return "\n".join(partes) if partes else None
    except ImportError:
        return None


def _extraer_xlsx(ruta: Path) -> Optional[str]:
    """Extrae texto de un archivo Excel (.xlsx)."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(ruta), data_only=True)
        partes: list[str] = []
        for hoja in wb.worksheets:
            for fila in hoja.iter_rows(values_only=True):
                celdas = [str(celda) for celda in fila if celda is not None and str(celda).strip()]
                if celdas:
                    partes.append(" ".join(celdas))
        return "\n".join(partes) if partes else None
    except ImportError:
        return None


def _extraer_pptx(ruta: Path) -> Optional[str]:
    """Extrae texto de un archivo PowerPoint (.pptx)."""
    try:
        from pptx import Presentation
        prs = Presentation(str(ruta))
        partes: list[str] = []
        for diapositiva in prs.slides:
            for forma in diapositiva.shapes:
                if hasattr(forma, "text") and forma.text.strip():
                    partes.append(forma.text.strip())
        return "\n".join(partes) if partes else None
    except ImportError:
        return None


def _extraer_imagen(ruta: Path) -> Optional[str]:
    """Extrae texto de una imagen usando Tesseract OCR."""
    try:
        from PIL import Image
        import pytesseract
        import os
        
        # Configurar la ruta de Tesseract en Windows si existe
        ruta_tesseract = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        if os.path.exists(ruta_tesseract):
            pytesseract.pytesseract.tesseract_cmd = ruta_tesseract
        
        imagen = Image.open(str(ruta))
        texto = pytesseract.image_to_string(imagen)
        return texto.strip() if texto.strip() else None
    except ImportError:
        return None
    except Exception as e:
        print(f"Error OCR en imagen {ruta}: {e}")
        return None


def _extraer_audio(ruta: Path) -> Optional[str]:
    """Extrae transcripción de un archivo de audio usando SpeechRecognition."""
    try:
        import speech_recognition as sr
        from moviepy.editor import AudioFileClip
        import tempfile
        import os
        
        # Convertir a WAV usando moviepy (que ya usa imageio_ffmpeg)
        audio = AudioFileClip(str(ruta))
        
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_wav:
            temp_wav_path = temp_wav.name
            
        # Limitar a los primeros 3 minutos (180 segundos) si es muy largo
        if audio.duration > 180:
            audio = audio.subclip(0, 180)
            
        audio.write_audiofile(temp_wav_path, logger=None)
        audio.close()
        
        recognizer = sr.Recognizer()
        with sr.AudioFile(temp_wav_path) as source:
            # Leer los primeros 3 minutos para no exceder límites de API
            audio_data = recognizer.record(source, duration=180)
            
        # Intentar reconocer (requiere internet para la API de Google)
        try:
            texto = recognizer.recognize_google(audio_data, language="es-ES")
        except sr.UnknownValueError:
            texto = None
        except sr.RequestError as e:
            print(f"Error en API de SpeechRecognition: {e}")
            texto = None
            
        # Limpiar temporal
        try:
            os.remove(temp_wav_path)
        except OSError:
            pass
            
        return texto.strip() if texto and texto.strip() else None
    except ImportError:
        return None
    except Exception as e:
        print(f"Error transcribiendo audio {ruta}: {e}")
        return None


def _extraer_video(ruta: Path) -> Optional[str]:
    """Extrae audio de un video y luego lo transcribe."""
    try:
        from moviepy.editor import VideoFileClip
        import tempfile
        import os
        
        video = VideoFileClip(str(ruta))
        if not video.audio:
            video.close()
            return None
            
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_audio:
            temp_audio_path = temp_audio.name
            
        # Extraer audio
        video.audio.write_audiofile(temp_audio_path, logger=None)
        video.close()
        
        # Reutilizar la función de extracción de audio
        texto = _extraer_audio(Path(temp_audio_path))
        
        # Limpiar temporal
        try:
            os.remove(temp_audio_path)
        except OSError:
            pass
            
        return texto
    except ImportError:
        return None
    except Exception as e:
        print(f"Error procesando video {ruta}: {e}")
        return None


def _extraer_texto_plano(ruta: Path) -> Optional[str]:
    """Lee archivos de texto plano con límite de tamaño."""
    try:
        tamaño = ruta.stat().st_size
        if tamaño > _MAX_PLAIN_BYTES:
            # Leer solo los primeros _MAX_PLAIN_BYTES
            with open(ruta, "r", encoding="utf-8", errors="replace") as f:
                return f.read(_MAX_PLAIN_BYTES)
        with open(ruta, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    except Exception:
        return None


def tokenizar_para_busqueda(texto: str | None) -> list[str]:
    """
    Convierte el texto en una lista de tokens únicos en minúsculas,
    filtrados por longitud mínima de 3 caracteres.
    Útil para generar palabras_clave automáticamente desde el contenido.
    """
    if not texto:
        return []
    tokens = _TOKEN_RE.findall(texto.lower())
    # Filtrar palabras muy cortas y stopwords básicas en español e inglés
    stopwords = {
        "de", "la", "el", "en", "y", "a", "los", "del", "se", "las", "por",
        "un", "para", "con", "una", "que", "su", "al", "es", "como", "si",
        "the", "and", "for", "are", "was", "that", "this", "with", "have",
        "from", "not", "but", "they", "can", "been", "has", "had",
    }
    vistos: set[str] = set()
    resultado: list[str] = []
    for token in tokens:
        if len(token) >= 3 and token not in stopwords and token not in vistos:
            vistos.add(token)
            resultado.append(token)
    return resultado
